import os
import json
import datetime
import time
import uuid
from typing import List, Optional
from fastapi import FastAPI, Depends, HTTPException, UploadFile, File, Form, Query, BackgroundTasks
from fastapi.staticfiles import StaticFiles
import shutil
from fastapi.middleware.cors import CORSMiddleware
from sqlalchemy.orm import Session
import pypdf
import docx2txt
from pptx import Presentation

from .database import engine, Base, get_db, SessionLocal
from . import models, schemas
from .auth import hash_password, verify_password, create_token, get_current_user
from .services.agents import (
    summarizer_agent,
    planner_recommender_agent,
    quiz_agent,
    tutor_agent,
    mock_exam_agent,
    formula_extractor_agent,
    deep_research_agent,
    curriculum_mapper_agent
)
from .services.vector_store import vector_store

# Create database tables
Base.metadata.create_all(bind=engine)

# Auto-migration: safely add material_id to flashcards table if it doesn't exist
try:
    from sqlalchemy import text
    with engine.begin() as conn:
        conn.execute(text("ALTER TABLE flashcards ADD COLUMN material_id INTEGER"))
except Exception as e:
    # Column already exists or table is not sqlite
    print(f"Migration note (can ignore if column already exists): {e}")

# Auto-migration: safely add deep_research_summary to materials table if it doesn't exist
try:
    from sqlalchemy import text
    with engine.begin() as conn:
        conn.execute(text("ALTER TABLE materials ADD COLUMN deep_research_summary TEXT"))
except Exception as e:
    print(f"Migration note (deep_research_summary already exists): {e}")

# Auto-migration: safely add explanation to quizzes table if it doesn't exist
try:
    from sqlalchemy import text
    with engine.begin() as conn:
        conn.execute(text("ALTER TABLE quizzes ADD COLUMN explanation TEXT"))
except Exception as e:
    print(f"Migration note (explanation already exists): {e}")

# Auto-migration: safely add title to study_sessions table if it doesn't exist
try:
    from sqlalchemy import text
    with engine.begin() as conn:
        conn.execute(text("ALTER TABLE study_sessions ADD COLUMN title VARCHAR"))
except Exception as e:
    print(f"Migration note (title already exists): {e}")

# Auto-migration: multi-tenancy — user_id on subjects (legacy rows keep NULL and
# become invisible to all accounts; users own only their own data from here on)
try:
    from sqlalchemy import text
    with engine.begin() as conn:
        conn.execute(text("ALTER TABLE subjects ADD COLUMN user_id INTEGER"))
except Exception as e:
    print(f"Migration note (user_id already exists): {e}")

# Per-user demo seeding: every new account starts with two sample subjects so
# the app is explorable immediately. Called from /api/auth/signup.
def seed_demo_data(db: Session, user_id: int):
    print(f"Seeding demo subjects for user {user_id}...")

    # 1. Subject 1: Operating Systems
    os_subj = models.Subject(
        user_id=user_id,
        name="Operating Systems (CS 401)",
        exam_date=(datetime.date.today() + datetime.timedelta(days=14)).isoformat(),
        priority_level=5,
        difficulty=4,
        confidence_score=45.0
    )
    db.add(os_subj)
    
    # 2. Subject 2: Computer Architecture
    ca_subj = models.Subject(
        user_id=user_id,
        name="Computer Architecture (CS 302)",
        exam_date=(datetime.date.today() + datetime.timedelta(days=3)).isoformat(),
        priority_level=4,
        difficulty=5,
        confidence_score=75.0
    )
    db.add(ca_subj)
    db.commit()
    db.refresh(os_subj)
    db.refresh(ca_subj)
    
    # Seed Tasks
    db.add_all([
        models.Task(
            subject_id=os_subj.id,
            title="Review Virtual Memory & Paging",
            description="Understand page tables, TLB hits/misses, and page replacement algorithms (LRU, FIFO).",
            duration_minutes=60,
            importance_score=9.0,
            urgency_score=7.0,
            status="pending",
            due_date=(datetime.date.today() + datetime.timedelta(days=4)).isoformat()
        ),
        models.Task(
            subject_id=os_subj.id,
            title="Study Process Synchronization & Semaphores",
            description="Resolve classical synchronization problems: Producer-Consumer, Readers-Writers.",
            duration_minutes=45,
            importance_score=8.5,
            urgency_score=6.0,
            status="pending",
            due_date=(datetime.date.today() + datetime.timedelta(days=6)).isoformat()
        ),
        models.Task(
            subject_id=os_subj.id,
            title="Syllabus Review & Key Concepts",
            description="Establish core exam topics, structural sections, and highlight highly rated exam concepts.",
            duration_minutes=30,
            importance_score=5.0,
            urgency_score=3.0,
            status="completed",
            due_date=datetime.date.today().isoformat(),
            completed_at=datetime.datetime.utcnow()
        ),
        models.Task(
            subject_id=ca_subj.id,
            title="Resolve Cache Mapping Exercises",
            description="Practice Direct, Fully Associative, and Set-Associative mapping cache addresses.",
            duration_minutes=90,
            importance_score=9.5,
            urgency_score=9.5,
            status="pending",
            due_date=(datetime.date.today() + datetime.timedelta(days=1)).isoformat()
        ),
        models.Task(
            subject_id=ca_subj.id,
            title="Review Instruction Pipelining Hazards",
            description="Differentiate Structural, Data, and Control hazards. Understand forwarding and branch prediction.",
            duration_minutes=50,
            importance_score=8.0,
            urgency_score=8.0,
            status="pending",
            due_date=(datetime.date.today() + datetime.timedelta(days=2)).isoformat()
        )
    ])
    
    # Seed Flashcards
    db.add_all([
        models.Flashcard(
            subject_id=os_subj.id,
            front="What is a Page Fault?",
            back="An interrupt raised by the hardware (MMU) when a program accesses a page that is mapped in the virtual address space, but not loaded in physical memory (RAM).",
            box=1,
            next_review_date=datetime.date.today().isoformat(),
            confidence=0.0
        ),
        models.Flashcard(
            subject_id=os_subj.id,
            front="Explain Thrashing",
            back="A state where the system spends more time swapping pages in and out of disk than executing instruction processes, caused by insufficient physical memory.",
            box=2,
            next_review_date=(datetime.date.today() + datetime.timedelta(days=2)).isoformat(),
            confidence=30.0
        )
    ])
    
    # Seed Quizzes
    db.add_all([
        models.Quiz(
            subject_id=os_subj.id,
            question="Which page replacement algorithm suffers from Belady's Anomaly?",
            options=json.dumps(["LRU (Least Recently Used)", "FIFO (First In First Out)", "Optimal", "Clock Replacement"]),
            correct_answer="FIFO (First In First Out)",
            type="multiple_choice"
        ),
        models.Quiz(
            subject_id=os_subj.id,
            question="What is the main advantage of a translation lookaside buffer (TLB)?",
            options=json.dumps(["Increases cache size", "Bypasses virtual memory", "Speeds up virtual-to-physical address translation", "Guarantees zero page faults"]),
            correct_answer="Speeds up virtual-to-physical address translation",
            type="multiple_choice"
        ),
        models.Quiz(
            subject_id=ca_subj.id,
            question="In instruction pipelining, which hazard is caused by memory or execution unit conflicts?",
            options=json.dumps(["Data hazard", "Control hazard", "Structural hazard", "Branch hazard"]),
            correct_answer="Structural hazard",
            type="multiple_choice"
        )
    ])
    db.commit()

app = FastAPI(title="Finals Buddy API", description="AI-powered University Finals Study Assistant Backend")

# Ensure uploads directory exists
os.makedirs("uploads", exist_ok=True)
app.mount("/uploads", StaticFiles(directory="uploads"), name="uploads")

# Setup CORS for frontend
# Local dev origins are always allowed; deployed frontend origins come from the
# CORS_ORIGINS env var (comma-separated), e.g. "https://finals-buddy.vercel.app"
from .config import CORS_ORIGINS

app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:3000",
        "http://127.0.0.1:3000",
        "http://localhost:3001",
        "http://127.0.0.1:3001",
        "http://localhost:3002",
        "http://127.0.0.1:3002",
        "http://localhost:3003",
        "http://127.0.0.1:3003",
        *CORS_ORIGINS,
    ],
    allow_origin_regex="https?://(localhost|127\\.0\\.0\\.1)(:\\d+)?",
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Helper to chunk text
def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start += chunk_size - overlap
    return chunks

# ----------------- AUTH -----------------

@app.post("/api/auth/signup", response_model=schemas.AuthResponse)
def signup(payload: schemas.SignupRequest, db: Session = Depends(get_db)):
    email = payload.email.strip().lower()
    if "@" not in email or "." not in email.split("@")[-1]:
        raise HTTPException(status_code=400, detail="Please enter a valid email address")
    if db.query(models.User).filter(models.User.email == email).first():
        raise HTTPException(status_code=409, detail="An account with this email already exists")

    user = models.User(email=email, name=payload.name.strip(), hashed_password=hash_password(payload.password))
    db.add(user)
    db.commit()
    db.refresh(user)

    # New accounts start empty — the dashboard shows a getting-started guide
    # until the user adds their first subject.
    return schemas.AuthResponse(token=create_token(user.id), user=user)

@app.post("/api/auth/login", response_model=schemas.AuthResponse)
def login(payload: schemas.LoginRequest, db: Session = Depends(get_db)):
    email = payload.email.strip().lower()
    user = db.query(models.User).filter(models.User.email == email).first()
    if not user or not verify_password(payload.password, user.hashed_password):
        raise HTTPException(status_code=401, detail="Incorrect email or password")
    return schemas.AuthResponse(token=create_token(user.id), user=user)

@app.get("/api/auth/me", response_model=schemas.UserOut)
def get_me(current_user: models.User = Depends(get_current_user)):
    return current_user

# ----------------- OWNERSHIP HELPERS -----------------

def own_subject(db: Session, user: models.User, subject_id: int) -> models.Subject:
    """Fetch a subject only if it belongs to the current user; 404 otherwise
    (404, not 403, so account enumeration isn't possible)."""
    subject = db.query(models.Subject).filter(
        models.Subject.id == subject_id,
        models.Subject.user_id == user.id
    ).first()
    if not subject:
        raise HTTPException(status_code=404, detail="Subject not found")
    return subject

def assert_owner(db: Session, user: models.User, subject_id: int):
    own_subject(db, user, subject_id)

# ----------------- SUBJECTS -----------------

@app.post("/api/subjects", response_model=schemas.SubjectOut)
def create_subject(subject: schemas.SubjectCreate, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    db_subj = models.Subject(**subject.model_dump(), user_id=current_user.id)
    db.add(db_subj)
    db.commit()
    db.refresh(db_subj)
    
    # Proactively generate a default first study plan task
    default_task = models.Task(
        subject_id=db_subj.id,
        title=f"Review syllabus & outline {db_subj.name}",
        description="Establish core exam topics, structural sections, and highlight highly rated exam concepts.",
        duration_minutes=45,
        importance_score=8.0,
        urgency_score=5.0,
        due_date=db_subj.exam_date or datetime.date.today().isoformat()
    )
    db.add(default_task)
    db.commit()
    return db_subj

@app.get("/api/subjects", response_model=List[schemas.SubjectOut])
def list_subjects(db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    return db.query(models.Subject).filter(models.Subject.user_id == current_user.id).all()

@app.get("/api/subjects/{subject_id}", response_model=schemas.SubjectDashboardOut)
def get_subject_detail(subject_id: int, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    subject = own_subject(db, current_user, subject_id)

    materials_count = len(subject.materials)
    
    # completion percentage
    total_tasks = len(subject.tasks)
    completed_tasks = len([t for t in subject.tasks if t.status == "completed"])
    completion_percentage = (completed_tasks / total_tasks * 100.0) if total_tasks > 0 else 0.0
    
    # hours remaining estimation
    pending_tasks = [t for t in subject.tasks if t.status == "pending"]
    hours_remaining = sum(t.duration_minutes for t in pending_tasks) / 60.0
    
    # Recalculate dynamic recommendations on the fly!
    recs = planner_recommender_agent.calculate_recommendations(subject, pending_tasks)
    
    # Weak topics scanner - aggregates from materials, recent low confidence
    weak_topics = []
    if subject.confidence_score < 60:
        weak_topics.append("Core Concepts Reinforcement")
    for mat in subject.materials:
        if mat.learning_complexity >= 4:
            weak_topics.append(mat.name)
            
    # Calculate Urgency Status based on exam date
    urgency = "low"
    if subject.exam_date:
        try:
            exam_dt = datetime.datetime.strptime(subject.exam_date, "%Y-%m-%d")
            days_left = (exam_dt - datetime.datetime.now()).days
            if days_left <= 2:
                urgency = "critical"
            elif days_left <= 7:
                urgency = "high"
            elif days_left <= 14:
                urgency = "medium"
        except Exception:
            pass
            
    next_action = recs[0]["reason"] if recs else "Add study materials or log a revision task to begin."
    if recs:
        # Get task details
        best_task = db.query(models.Task).filter(models.Task.id == recs[0]["task_id"]).first()
        if best_task:
            next_action = f"Study '{best_task.title}' ({recs[0]['score']}% Priority Match)"

    return schemas.SubjectDashboardOut(
        id=subject.id,
        name=subject.name,
        exam_date=subject.exam_date,
        priority_level=subject.priority_level,
        difficulty=subject.difficulty,
        confidence_score=subject.confidence_score,
        materials_count=materials_count,
        completion_percentage=round(completion_percentage, 1),
        hours_remaining=round(hours_remaining, 1),
        weak_topics=weak_topics[:4] or ["Ready for Testing"],
        next_recommended_action=next_action,
        urgency_status=urgency
    )

@app.patch("/api/subjects/{subject_id}", response_model=schemas.SubjectOut)
def update_subject(subject_id: int, subject_update: schemas.SubjectUpdate, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    subject = own_subject(db, current_user, subject_id)

    update_data = subject_update.model_dump(exclude_unset=True)
    for key, value in update_data.items():
        setattr(subject, key, value)

    db.commit()
    db.refresh(subject)
    return subject

@app.delete("/api/subjects/{subject_id}")
def delete_subject(subject_id: int, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    subject = own_subject(db, current_user, subject_id)
    db.delete(subject)
    db.commit()
    return {"message": "Subject deleted successfully"}

# Global in-memory dictionary to track background ingestion progress
ingestion_jobs = {}

class IngestionProgressTracker:
    def __init__(self, job_id: str):
        self.job_id = job_id
        if job_id not in ingestion_jobs:
            ingestion_jobs[job_id] = {
                "step": 1,
                "label": "Preparing...",
                "icon": "📤",
                "status": "processing",
                "timestamp": datetime.datetime.now().isoformat()
            }
        
    def update(self, step: int, label: str, icon: str, status: str = "processing", data: dict = None):
        if self.job_id not in ingestion_jobs:
            ingestion_jobs[self.job_id] = {}
        payload = {
            "step": step,
            "label": label,
            "icon": icon,
            "status": status,
            "timestamp": datetime.datetime.now().isoformat()
        }
        if data:
            payload.update(data)
        ingestion_jobs[self.job_id].update(payload)


def process_ingestion_background(
    job_id: str,
    material_id: int,
    subject_id: int,
    filename: str,
    file_ext: str,
    file_path: str,
    text_content: str
):
    tracker = IngestionProgressTracker(job_id)
    db = SessionLocal()
    t_start = time.time()
    try:
        # Step 3: Chunk & Index
        tracker.update(3, "Chunking & indexing for RAG search...", "🗂️")
        t2 = time.time()
        chunks = chunk_text(text_content)
        print(f"🔪  [3/7] Background: Chunking → {len(chunks)} chunks...")
        for idx, chunk in enumerate(chunks):
            labeled_chunk = f"[Source Document: {filename}]\n{chunk}"
            vector_store.add_document(
                text=labeled_chunk,
                metadata={"subject_id": subject_id, "name": filename, "chunk_index": idx}
            )
        print(f"🗂️   [3/7] Background: Indexed in {time.time()-t2:.2f}s")
        
        # Step 4: AI Summarization
        tracker.update(4, "AI generating deep summary with citations...", "🧠")
        t3 = time.time()
        print("🧠  [4/7] Background: Running AI Summarization Agent...")
        analysis = summarizer_agent.process_material(filename, text_content)
        print(f"✅  [4/7] Background: Summarization complete ({time.time()-t3:.2f}s)")
        
        # Step 5: Deep Research
        tracker.update(5, "Deep research enrichment pass...", "🔬")
        t4 = time.time()
        print("🔬  [5/7] Background: Running Deep Research Agent...")
        deep_summary = deep_research_agent.enrich_material(filename, text_content)
        print(f"✅  [5/7] Background: Deep research complete ({time.time()-t4:.2f}s)")
        
        # Save steps 1-5 results to Material DB
        t5 = time.time()
        db_material = db.query(models.Material).filter(models.Material.id == material_id).first()
        if db_material:
            db_material.summary = analysis.get("summary", "")
            db_material.key_concepts = json.dumps(analysis.get("key_concepts", []))
            db_material.learning_complexity = analysis.get("learning_complexity", 3)
            db_material.importance_level = analysis.get("importance_level", 4)
            db_material.deep_research_summary = deep_summary
            db.commit()
            db.refresh(db_material)
            print(f"💽  [5/7] Background: Material saved to DB ({time.time()-t5:.2f}s)")
            
        # We notify progress stream that the core material summary is fully ready and saved!
        # This allows the frontend to immediately unlock the dashboard for reading!
        tracker.update(5, "Summary ready! Generating flashcards...", "📝", data={
            "summary_ready": True,
            "material_id": material_id
        })
        
        # Step 6: Quiz & Flashcard generation
        tracker.update(6, "Generating cited flashcards & quiz questions...", "📝")
        t6 = time.time()
        cited_summary = analysis.get("summary", "")
        print("📝  [6/7] Background: Running Quiz & Flashcard Agent...")
        active_recall = quiz_agent.generate_quiz_and_flashcards(filename, text_content, cited_summary)
        
        flashcard_count = len(active_recall.get("flashcards", []))
        quiz_count = len(active_recall.get("quizzes", []))
        
        for fc in active_recall.get("flashcards", []):
            db_fc = models.Flashcard(
                subject_id=subject_id,
                material_id=material_id,
                front=fc["front"],
                back=fc["back"],
                box=1,
                next_review_date=datetime.date.today().isoformat()
            )
            db.add(db_fc)
            
        for q in active_recall.get("quizzes", []):
            db_q = models.Quiz(
                subject_id=subject_id,
                material_id=material_id,
                question=q["question"],
                options=json.dumps(q["options"]),
                correct_answer=q["correct_answer"],
                explanation=q.get("explanation", ""),
                type="multiple_choice"
            )
            db.add(db_q)
        db.commit()
        print(f"✅  [6/7] Background: Active recall generated ({time.time()-t6:.2f}s)")
        
        # Step 7: Create dynamic study task
        tracker.update(7, "Finalizing & creating study tasks...", "📋")
        t7 = time.time()
        study_task = models.Task(
            subject_id=subject_id,
            title=f"Review summary & active recall for {filename}",
            description=f"Go through the dynamic summary and test your recall on the key concepts derived from {filename}.",
            duration_minutes=30,
            importance_score=float(analysis.get("importance_level", 4) * 2),
            urgency_score=6.0,
            due_date=(datetime.date.today() + datetime.timedelta(days=2)).isoformat()
        )
        db.add(study_task)
        db.commit()
        print(f"📋  [7/7] Background: Study task created ({time.time()-t7:.2f}s)")
        
        total = time.time() - t_start
        print(f"\n{'='*60}")
        print(f"🎉  [INGESTION COMPLETE] '{filename}' processed in {total:.1f}s")
        print(f"     Summary: {len(analysis.get('summary',''))} chars | Flashcards: {flashcard_count} | Quizzes: {quiz_count}")
        print(f"{'='*60}\n")
        
        # Mark complete!
        tracker.update(7, "Complete! Loading your new study materials...", "🎉", status="completed", data={
            "summary_ready": True,
            "material_id": material_id
        })
    except Exception as e:
        print(f"❌ Background digestion error: {e}")
        import traceback
        traceback.print_exc()
        tracker.update(7, f"Digestion failed: {str(e)}", "❌", status="failed", data={"error": str(e)})
    finally:
        db.close()


@app.post("/api/materials/upload", response_model=schemas.MaterialOut)
async def upload_material(
    background_tasks: BackgroundTasks,
    subject_id: int = Form(...),
    file: UploadFile = File(...),
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user)
):
    t_start = time.time()

    subject = own_subject(db, current_user, subject_id)

    filename = file.filename
    file_ext = filename.split(".")[-1].lower()
    file_bytes = await file.read()
    file_size_kb = len(file_bytes) / 1024

    print(f"\n{'='*60}")
    print(f"📥  [INGESTION START] '{filename}' ({file_size_kb:.1f} KB) → Subject: '{subject.name}' (id={subject_id})")
    print(f"{'='*60}")

    # Save the file locally in a workspace upload directory
    upload_dir = "./uploads"
    os.makedirs(upload_dir, exist_ok=True)
    file_path = os.path.join(upload_dir, f"{subject_id}_{int(datetime.datetime.now().timestamp())}_{filename}")

    with open(file_path, "wb") as f:
        f.write(file_bytes)
    print(f"💾  [1/7] File saved → {file_path}")

    # Extract text from file
    text_content = ""
    t1 = time.time()
    try:
        if file_ext == "pdf":
            reader = pypdf.PdfReader(file_path)
            for page in reader.pages:
                text_content += page.extract_text() or ""
            print(f"📄  [2/7] PDF text extracted — {len(reader.pages)} pages, {len(text_content):,} chars  ({time.time()-t1:.2f}s)")
        elif file_ext in ["docx", "doc"]:
            text_content = docx2txt.process(file_path)
            print(f"📄  [2/7] DOCX text extracted — {len(text_content):,} chars  ({time.time()-t1:.2f}s)")
        elif file_ext in ["pptx", "ppt"]:
            prs = Presentation(file_path)
            slide_texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
            text_content = "\n".join(slide_texts)
            print(f"📄  [2/7] PPTX text extracted — {len(prs.slides)} slides, {len(text_content):,} chars  ({time.time()-t1:.2f}s)")
        else:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as tf:
                text_content = tf.read()
            print(f"📄  [2/7] Plain text extracted — {len(text_content):,} chars  ({time.time()-t1:.2f}s)")
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to parse uploaded document: {str(e)}")

    if not text_content.strip():
        raise HTTPException(status_code=400, detail="Document text content appears empty.")

    # Save initial placeholder material in relational DB immediately so we can get its ID
    db_material = models.Material(
        subject_id=subject_id,
        name=filename,
        file_type=file_ext,
        file_path=file_path,
        summary="",
        key_concepts="[]",
        learning_complexity=3,
        importance_level=4,
        deep_research_summary=""
    )
    db.add(db_material)
    db.commit()
    db.refresh(db_material)
    print(f"💽  [INITIAL] Material placeholder saved to DB — id={db_material.id}")

    # Generate a unique job ID for progress streaming
    job_id = f"job_{db_material.id}_{uuid.uuid4().hex[:8]}"
    
    # Initialize the progress tracker
    tracker = IngestionProgressTracker(job_id)
    tracker.update(1, "File saved locally...", "💾", status="processing")
    tracker.update(2, "Text extracted successfully...", "📄", status="processing")

    # Add ingestion background task to FastAPI's non-blocking background queue
    background_tasks.add_task(
        process_ingestion_background,
        job_id=job_id,
        material_id=db_material.id,
        subject_id=subject_id,
        filename=filename,
        file_ext=file_ext,
        file_path=file_path,
        text_content=text_content
    )
    
    # Inject job_id in dynamic attribute so it matches schemas.MaterialOut
    db_material.job_id = job_id
    
    print(f"🚀  Background digestion queued! Job ID: {job_id}")
    return db_material


@app.get("/api/materials/upload-progress/{job_id}")
async def upload_progress_stream(job_id: str):
    """
    Server-Sent Events (SSE) endpoint to stream document digestion logs & progress in real-time.
    """
    from fastapi.responses import StreamingResponse
    import asyncio

    async def event_generator():
        print(f"📡 Progress Stream: Client connected to job: {job_id}")
        last_step = -1
        
        while True:
            job = ingestion_jobs.get(job_id)
            if not job:
                await asyncio.sleep(0.5)
                continue
                
            # Stream the updated progress
            yield f"data: {json.dumps(job)}\n\n"
            
            if job.get("status") in ["completed", "failed"]:
                print(f"📡 Progress Stream: Completed for job: {job_id}")
                break
                
            await asyncio.sleep(1.0)
            
    return StreamingResponse(event_generator(), media_type="text/event-stream")



@app.get("/api/subjects/{subject_id}/materials", response_model=List[schemas.MaterialOut])
def list_subject_materials(subject_id: int, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    assert_owner(db, current_user, subject_id)
    return db.query(models.Material).filter(models.Material.subject_id == subject_id).all()

@app.post("/api/subjects/{subject_id}/generate-map", response_model=schemas.KnowledgeMapOut)
def generate_knowledge_map(subject_id: int, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    assert_owner(db, current_user, subject_id)
    # 1. Fetch all materials for this subject
    materials = db.query(models.Material).filter(models.Material.subject_id == subject_id).all()
    if len(materials) < 2:
        # Not enough materials to map, clear and return empty edges
        db.query(models.ResourceConnection).filter(models.ResourceConnection.subject_id == subject_id).delete()
        db.commit()
        return {"nodes": materials, "edges": []}
        
    # 2. Compile list of material summaries to send to CurriculumMapperAgent
    materials_info = []
    for m in materials:
        materials_info.append({
            "name": m.name,
            "summary": m.summary or ""
        })
        
    # 3. Trigger CurriculumMapperAgent
    map_result = curriculum_mapper_agent.generate_material_map(materials_info)
    
    # 4. Prune existing connections for this subject
    db.query(models.ResourceConnection).filter(models.ResourceConnection.subject_id == subject_id).delete()
    db.commit()
    
    # 5. Insert newly compiled connections
    edges = []
    for conn in map_result.get("connections", []):
        source = db.query(models.Material).filter(
            models.Material.subject_id == subject_id,
            models.Material.name == conn["source_material_name"]
        ).first()
        target = db.query(models.Material).filter(
            models.Material.subject_id == subject_id,
            models.Material.name == conn["target_material_name"]
        ).first()
        
        # Resilient mapping: only insert if source and target actually match ingested materials
        if source and target:
            db_conn = models.ResourceConnection(
                subject_id=subject_id,
                source_material_id=source.id,
                target_material_id=target.id,
                connection_type=conn["connection_type"],
                description=conn.get("description", "")
            )
            db.add(db_conn)
            edges.append(db_conn)
            
    db.commit()
    
    # Refresh edges
    for e in edges:
        db.refresh(e)
        
    return {"nodes": materials, "edges": edges}

@app.get("/api/subjects/{subject_id}/map", response_model=schemas.KnowledgeMapOut)
def get_knowledge_map(subject_id: int, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    assert_owner(db, current_user, subject_id)
    materials = db.query(models.Material).filter(models.Material.subject_id == subject_id).all()
    edges = db.query(models.ResourceConnection).filter(models.ResourceConnection.subject_id == subject_id).all()
    return {"nodes": materials, "edges": edges}

@app.patch("/api/materials/{material_id}", response_model=schemas.MaterialOut)
def update_material(material_id: int, material_update: schemas.MaterialUpdate, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    material = db.query(models.Material).filter(models.Material.id == material_id).first()
    if not material:
        raise HTTPException(status_code=404, detail="Material not found")
    assert_owner(db, current_user, material.subject_id)

    update_data = material_update.model_dump(exclude_unset=True)
    for key, value in update_data.items():
        setattr(material, key, value)
        
    db.commit()
    db.refresh(material)
    return material

@app.delete("/api/materials/{material_id}")
def delete_material(material_id: int, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    material = db.query(models.Material).filter(models.Material.id == material_id).first()
    if not material:
        raise HTTPException(status_code=404, detail="Material not found")
    assert_owner(db, current_user, material.subject_id)

    # 1. Remove physical file from disk
    if material.file_path and os.path.exists(material.file_path):
        try:
            os.remove(material.file_path)
        except Exception as e:
            print(f"Error removing physical file: {e}")
            
    # 2. Clean up indexed document chunks from vector store
    try:
        vector_store.remove_documents_by_metadata({
            "subject_id": material.subject_id,
            "name": material.name
        })
    except Exception as e:
        print(f"Error removing vector store index: {e}")
        
    # 3. Delete from database
    db.delete(material)
    db.commit()
    return {"message": "Material and associated indexes deleted successfully"}

# ----------------- ADAPTIVE STUDY PLANNER & TASKS -----------------

@app.post("/api/tasks", response_model=schemas.TaskOut)
def create_task(task: schemas.TaskCreate, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    assert_owner(db, current_user, task.subject_id)
    db_task = models.Task(**task.model_dump())
    db.add(db_task)
    db.commit()
    db.refresh(db_task)
    return db_task

@app.get("/api/subjects/{subject_id}/tasks", response_model=List[schemas.TaskOut])
def list_subject_tasks(subject_id: int, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    assert_owner(db, current_user, subject_id)
    return db.query(models.Task).filter(models.Task.subject_id == subject_id).all()

@app.delete("/api/tasks/{task_id}")
def delete_task(task_id: int, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    db_task = db.query(models.Task).filter(models.Task.id == task_id).first()
    if not db_task:
        raise HTTPException(status_code=404, detail="Task not found")
    assert_owner(db, current_user, db_task.subject_id)
    db.delete(db_task)
    db.commit()
    return {"message": "Task deleted successfully"}

@app.patch("/api/tasks/{task_id}", response_model=schemas.TaskOut)
def update_task(task_id: int, task_data: schemas.TaskUpdate, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    db_task = db.query(models.Task).filter(models.Task.id == task_id).first()
    if not db_task:
        raise HTTPException(status_code=404, detail="Task not found")
    assert_owner(db, current_user, db_task.subject_id)

    for key, value in task_data.model_dump(exclude_unset=True).items():
        setattr(db_task, key, value)
        
    if task_data.status == "completed":
        db_task.completed_at = datetime.datetime.utcnow()
        # Completing a task slightly boosts overall subject confidence score!
        subject = db_task.subject
        if subject:
            subject.confidence_score = min(100.0, subject.confidence_score + 3.5)
            
    db.commit()
    db.refresh(db_task)
    return db_task

@app.post("/api/study-sessions", response_model=schemas.StudySessionOut)
def create_study_session(session: schemas.StudySessionCreate, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    assert_owner(db, current_user, session.subject_id)
    db_session = models.StudySession(**session.model_dump())
    db.add(db_session)
    
    # Increment confidence slightly based on high focus score
    subject = db.query(models.Subject).filter(models.Subject.id == session.subject_id).first()
    if subject and session.focus_score:
        subject.confidence_score = min(100.0, subject.confidence_score + (session.focus_score * 0.8))

    db.commit()
    db.refresh(db_session)
    return db_session

@app.get("/api/subjects/{subject_id}/sessions", response_model=List[schemas.StudySessionOut])
def list_subject_sessions(subject_id: int, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    assert_owner(db, current_user, subject_id)
    return db.query(models.StudySession).filter(models.StudySession.subject_id == subject_id).all()

# ----------------- NOTES / SCRATCHPAD -----------------

@app.delete("/api/study-sessions/{session_id}")
def delete_study_session(session_id: int, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    db_session = db.query(models.StudySession).filter(models.StudySession.id == session_id).first()
    if not db_session:
        raise HTTPException(status_code=404, detail="Study session not found")
    assert_owner(db, current_user, db_session.subject_id)
    db.delete(db_session)
    db.commit()
    return {"message": "Study session deleted successfully"}

@app.get("/api/subjects/{subject_id}/notes", response_model=List[schemas.NoteOut])
def list_subject_notes(subject_id: int, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    assert_owner(db, current_user, subject_id)
    return db.query(models.Note).filter(models.Note.subject_id == subject_id).order_by(models.Note.updated_at.desc()).all()

@app.post("/api/subjects/{subject_id}/notes", response_model=schemas.NoteOut)
def create_subject_note(subject_id: int, note: schemas.NoteCreate, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    assert_owner(db, current_user, subject_id)
    db_note = models.Note(
        subject_id=subject_id,
        title=note.title,
        content=note.content
    )
    db.add(db_note)
    db.commit()
    db.refresh(db_note)
    return db_note

@app.patch("/api/notes/{note_id}", response_model=schemas.NoteOut)
def update_note(note_id: int, note_update: schemas.NoteUpdate, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    db_note = db.query(models.Note).filter(models.Note.id == note_id).first()
    if not db_note:
        raise HTTPException(status_code=404, detail="Note not found")
    assert_owner(db, current_user, db_note.subject_id)
    
    if note_update.title is not None:
        db_note.title = note_update.title
    if note_update.content is not None:
        db_note.content = note_update.content
        
    db.commit()
    db.refresh(db_note)
    return db_note

@app.delete("/api/notes/{note_id}")
def delete_note(note_id: int, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    db_note = db.query(models.Note).filter(models.Note.id == note_id).first()
    if not db_note:
        raise HTTPException(status_code=404, detail="Note not found")
    assert_owner(db, current_user, db_note.subject_id)
    db.delete(db_note)
    db.commit()
    return {"message": "Note deleted successfully"}

# ----------------- UPLOADS -----------------

@app.post("/api/upload")
async def upload_file(file: UploadFile = File(...), current_user: models.User = Depends(get_current_user)):
    # Generate unique filename to prevent overwrites
    ext = file.filename.split(".")[-1] if "." in file.filename else ""
    filename = f"{uuid.uuid4().hex}.{ext}"
    file_path = os.path.join("uploads", filename)
    
    with open(file_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
        
    return {"url": f"http://localhost:8000/uploads/{filename}"}

# ----------------- REVISION & LEITNER FLASHCARDS -----------------

@app.get("/api/subjects/{subject_id}/flashcards", response_model=List[schemas.FlashcardOut])
def list_subject_flashcards(subject_id: int, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    assert_owner(db, current_user, subject_id)
    # Return all flashcards so frontend can filter by due date or view all in manage mode
    return db.query(models.Flashcard).filter(models.Flashcard.subject_id == subject_id).all()

@app.post("/api/subjects/{subject_id}/flashcards", response_model=schemas.FlashcardOut)
def create_subject_flashcard(subject_id: int, flashcard: schemas.FlashcardCreate, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    assert_owner(db, current_user, subject_id)
    db_card = models.Flashcard(
        subject_id=subject_id,
        front=flashcard.front,
        back=flashcard.back,
        box=flashcard.box or 1,
        next_review_date=datetime.date.today().isoformat()
    )
    db.add(db_card)
    db.commit()
    db.refresh(db_card)
    return db_card

@app.delete("/api/flashcards/{card_id}")
def delete_flashcard(card_id: int, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    card = db.query(models.Flashcard).filter(models.Flashcard.id == card_id).first()
    if not card:
        raise HTTPException(status_code=404, detail="Flashcard not found")
    assert_owner(db, current_user, card.subject_id)
    db.delete(card)
    db.commit()
    return {"message": "Flashcard deleted successfully"}

@app.put("/api/flashcards/{card_id}", response_model=schemas.FlashcardOut)
def update_flashcard(card_id: int, updates: schemas.FlashcardUpdate, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    card = db.query(models.Flashcard).filter(models.Flashcard.id == card_id).first()
    if not card:
        raise HTTPException(status_code=404, detail="Flashcard not found")
    assert_owner(db, current_user, card.subject_id)
    
    if updates.front is not None:
        card.front = updates.front
    if updates.back is not None:
        card.back = updates.back
    if updates.box is not None:
        card.box = updates.box
    if updates.material_id is not None:
        card.material_id = updates.material_id
        
    db.commit()
    db.refresh(card)
    return card

@app.post("/api/flashcards/{card_id}/review", response_model=schemas.FlashcardOut)
def review_flashcard(card_id: int, review: schemas.FlashcardReviewRequest, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    card = db.query(models.Flashcard).filter(models.Flashcard.id == card_id).first()
    if not card:
        raise HTTPException(status_code=404, detail="Flashcard not found")
    assert_owner(db, current_user, card.subject_id)
        
    # Leitner Spaced Repetition Logic
    if review.is_correct:
        card.box = min(5, card.box + 1)
        card.confidence = min(100.0, card.confidence + 15.0)
    else:
        card.box = 1
        card.confidence = max(0.0, card.confidence - 25.0)
        
    # Determine next review date based on Leitner Box
    # Box 1: 1 day, Box 2: 3 days, Box 3: 7 days, Box 4: 14 days, Box 5: 30 days
    intervals = {1: 1, 2: 3, 3: 7, 4: 14, 5: 30}
    days = intervals.get(card.box, 1)
    
    card.next_review_date = (datetime.date.today() + datetime.timedelta(days=days)).isoformat()
    db.commit()
    db.refresh(card)
    return card

@app.post("/api/subjects/{subject_id}/generate-more")
def generate_more_active_recall(subject_id: int, request: schemas.GenerateMoreRequest, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    subject = own_subject(db, current_user, subject_id)

    materials = db.query(models.Material).filter(models.Material.subject_id == subject_id)
    if request.material_id is not None:
        materials = materials.filter(models.Material.id == request.material_id)
    materials = materials.all()

    combined_summary = "\n\n".join([m.summary for m in materials if m.summary])
    if not combined_summary:
        combined_summary = f"Generic knowledge about {subject.name}"

    existing_questions = []
    if request.item_type == 'flashcards':
        existing_cards = db.query(models.Flashcard).filter(models.Flashcard.subject_id == subject_id).all()
        existing_questions = [c.front for c in existing_cards]
    else:
        existing_quizzes = db.query(models.Quiz).filter(models.Quiz.subject_id == subject_id).all()
        existing_questions = [q.question for q in existing_quizzes]

    generated_data = quiz_agent.generate_more_items(
        context_summary=combined_summary,
        existing_questions=existing_questions,
        item_type=request.item_type,
        count=request.count or 3
    )

    new_items = []
    mat_id_to_use = request.material_id if request.material_id else (materials[-1].id if materials else None)

    if request.item_type == 'flashcards':
        for fc in generated_data.get('flashcards', []):
            db_fc = models.Flashcard(
                subject_id=subject_id,
                material_id=mat_id_to_use,
                front=fc.get("front", ""),
                back=fc.get("back", ""),
                box=1,
                next_review_date=datetime.date.today().isoformat()
            )
            db.add(db_fc)
            new_items.append(db_fc)
    else:
        for q in generated_data.get('quizzes', []):
            db_q = models.Quiz(
                subject_id=subject_id,
                material_id=mat_id_to_use,
                question=q.get("question", ""),
                options=json.dumps(q.get("options", [])),
                correct_answer=q.get("correct_answer", ""),
                explanation=q.get("explanation", ""),
                type="multiple_choice"
            )
            db.add(db_q)
            new_items.append(db_q)

    db.commit()
    # Need to return dicts or Pydantic models. We will convert SQL objects to dicts manually 
    # to avoid missing attributes, or just return the schema-compatible objects.
    result = []
    for item in new_items:
        db.refresh(item)
        result.append(item)
        
    return result

@app.get("/api/subjects/{subject_id}/quizzes", response_model=List[schemas.QuizOut])
def list_subject_quizzes(subject_id: int, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    assert_owner(db, current_user, subject_id)
    return db.query(models.Quiz).filter(models.Quiz.subject_id == subject_id).all()

@app.post("/api/subjects/{subject_id}/quizzes", response_model=schemas.QuizOut)
def create_quiz_question(subject_id: int, quiz: schemas.QuizBase, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    assert_owner(db, current_user, subject_id)
    db_quiz = models.Quiz(subject_id=subject_id, **quiz.model_dump())
    db.add(db_quiz)
    db.commit()
    db.refresh(db_quiz)
    return db_quiz

@app.put("/api/quizzes/{quiz_id}", response_model=schemas.QuizOut)
def update_quiz_question(quiz_id: int, updates: schemas.QuizUpdate, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    quiz = db.query(models.Quiz).filter(models.Quiz.id == quiz_id).first()
    if not quiz:
        raise HTTPException(status_code=404, detail="Quiz question not found")
    assert_owner(db, current_user, quiz.subject_id)
    for key, value in updates.model_dump(exclude_unset=True).items():
        setattr(quiz, key, value)
    db.commit()
    db.refresh(quiz)
    return quiz

@app.delete("/api/quizzes/{quiz_id}")
def delete_quiz_question(quiz_id: int, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    quiz = db.query(models.Quiz).filter(models.Quiz.id == quiz_id).first()
    if not quiz:
        raise HTTPException(status_code=404, detail="Quiz question not found")
    assert_owner(db, current_user, quiz.subject_id)
    db.delete(quiz)
    db.commit()
    return {"message": "Quiz question deleted successfully"}

@app.post("/api/quizzes/{quiz_id}/answer", response_model=schemas.QuizAnswerResponse)
def answer_quiz_question(quiz_id: int, answer_req: schemas.QuizAnswerRequest, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    quiz = db.query(models.Quiz).filter(models.Quiz.id == quiz_id).first()
    if not quiz:
        raise HTTPException(status_code=404, detail="Quiz question not found")
    assert_owner(db, current_user, quiz.subject_id)
        
    is_correct = quiz.correct_answer.strip().lower() == answer_req.user_answer.strip().lower()
    
    # Adjust subject confidence score
    subject = quiz.subject
    if subject:
        if is_correct:
            subject.confidence_score = min(100.0, subject.confidence_score + 2.0)
        else:
            subject.confidence_score = max(0.0, subject.confidence_score - 4.0)
        db.commit()
        
    # Use stored AI-generated cited explanation if available, else fallback
    ai_explanation = quiz.explanation or ""
    if not ai_explanation:
        ai_explanation = "Perfect! Spaced recall reinforcement registered." if is_correct else "Review this topic again in your dynamic planners."
    
    return schemas.QuizAnswerResponse(
        is_correct=is_correct,
        correct_answer=quiz.correct_answer,
        explanation=ai_explanation
    )

# ----------------- AI TUTOR (RAG) -----------------

@app.post("/api/tutor/chat")
def tutor_chat(
    subject_id: int = Form(...),
    query: str = Form(...),
    mode: str = Form("standard"),
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user)
):
    subject = own_subject(db, current_user, subject_id)
        
    # 1. Save user message
    user_msg = models.ChatMessage(
        subject_id=subject_id,
        role="user",
        content=query
    )
    db.add(user_msg)
    db.commit()
    
    # 2. Get AI tutor response via LangChain agent (local Ollama + tools + memory)
    from .services.langchain_chat import run_langchain_chat
    response = run_langchain_chat(subject_id, query, mode, db)
    
    # 3. Save assistant response with sources inline
    ans_content = response.get("answer", "")
    sources = response.get("sources", [])
    if sources:
        ans_content += f"\n\n[Sources: {', '.join(sources)}]"
        
    assistant_msg = models.ChatMessage(
        subject_id=subject_id,
        role="assistant",
        content=ans_content
    )
    db.add(assistant_msg)
    db.commit()
    
    return {
        "answer": response.get("answer", ""),
        "sources": sources,
        "user_message_id": user_msg.id,
        "assistant_message_id": assistant_msg.id
    }

@app.get("/api/subjects/{subject_id}/chats", response_model=List[schemas.ChatMessageOut])
def get_chat_history(subject_id: int, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    assert_owner(db, current_user, subject_id)
    return db.query(models.ChatMessage).filter(models.ChatMessage.subject_id == subject_id).order_by(models.ChatMessage.created_at.asc()).all()

@app.patch("/api/chats/{message_id}", response_model=schemas.ChatMessageOut)
def update_chat_message(message_id: int, message_update: schemas.ChatMessageUpdate, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    msg = db.query(models.ChatMessage).filter(models.ChatMessage.id == message_id).first()
    if not msg:
        raise HTTPException(status_code=404, detail="Message not found")
    assert_owner(db, current_user, msg.subject_id)
    msg.content = message_update.content
    db.commit()
    db.refresh(msg)
    return msg

@app.delete("/api/subjects/{subject_id}/chats")
def clear_chat_history(subject_id: int, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    assert_owner(db, current_user, subject_id)
    
    db.query(models.ChatMessage).filter(models.ChatMessage.subject_id == subject_id).delete()
    db.commit()
    return {"message": "Chat history cleared successfully"}


# ----------------- GLOBAL ANALYTICS & DASHBOARD -----------------

@app.get("/api/dashboard/recommendations", response_model=List[schemas.RecommendationOut])
def get_global_recommendations(db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    subjects = db.query(models.Subject).filter(models.Subject.user_id == current_user.id).all()
    all_recs = []
    
    for subject in subjects:
        pending_tasks = db.query(models.Task).filter(
            models.Task.subject_id == subject.id,
            models.Task.status == "pending"
        ).all()
        
        recs = planner_recommender_agent.calculate_recommendations(subject, pending_tasks)
        
        for r in recs:
            # Check if recommendation already exists to prevent duplicate insertion
            db_rec = db.query(models.Recommendation).filter(
                models.Recommendation.subject_id == subject.id,
                models.Recommendation.task_id == r["task_id"]
            ).first()
            
            if not db_rec:
                db_rec = models.Recommendation(
                    subject_id=subject.id,
                    task_id=r["task_id"],
                    score=r["score"],
                    reason=r["reason"]
                )
                db.add(db_rec)
                
            all_recs.append(db_rec)
            
    db.commit()
    
    # Return top active recommendations sorted by score (scoped to this user)
    user_subject_ids = [s.id for s in subjects]
    active_recs = db.query(models.Recommendation).filter(
        models.Recommendation.is_dismissed == False,
        models.Recommendation.subject_id.in_(user_subject_ids)
    ).order_by(models.Recommendation.score.desc()).all()
    
    return active_recs[:5]

@app.get("/api/dashboard/summary")
def get_dashboard_summary(db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    subjects = db.query(models.Subject).filter(models.Subject.user_id == current_user.id).all()
    total_subjects = len(subjects)
    subject_ids = [s.id for s in subjects]

    # Aggregations (scoped to this user)
    total_study_minutes = db.query(models.StudySession).filter(
        models.StudySession.subject_id.in_(subject_ids)
    ).with_entities(
        models.StudySession.duration_minutes
    ).all()
    studied_hours = sum([m[0] for m in total_study_minutes if m[0]]) / 60.0

    tasks = db.query(models.Task).filter(models.Task.subject_id.in_(subject_ids)).all()
    total_tasks = len(tasks)
    completed_tasks = len([t for t in tasks if t.status == "completed"])
    global_completion = (completed_tasks / total_tasks * 100.0) if total_tasks > 0 else 0.0
    
    # Calculate average confidence
    avg_confidence = sum([s.confidence_score for s in subjects]) / total_subjects if total_subjects > 0 else 0.0
    
    # Burnout index: High hours studied in last 3 days + late hours + high subject stress
    burnout_index = min(100.0, (studied_hours * 2.5) + (total_tasks - completed_tasks) * 1.5)
    
    # Mock data for SVGs
    subject_names = [s.name for s in subjects]
    confidence_data = [s.confidence_score for s in subjects]
    
    return {
        "total_subjects": total_subjects,
        "studied_hours": round(studied_hours, 1),
        "completion_rate": round(global_completion, 1),
        "average_confidence": round(avg_confidence, 1),
        "burnout_risk_percentage": round(burnout_index, 1),
        "subject_analytics": {
            "labels": subject_names,
            "values": confidence_data
        }
    }


# ----------------- PHASE 2: MOCK EXAMS & CHEAT SHEETS -----------------

@app.get("/api/subjects/{subject_id}/mock-exams", response_model=List[schemas.MockExamOut])
def list_subject_mock_exams(subject_id: int, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    assert_owner(db, current_user, subject_id)
    return db.query(models.MockExam).filter(models.MockExam.subject_id == subject_id).all()

@app.delete("/api/mock-exams/{exam_id}")
def delete_mock_exam(exam_id: int, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    exam = db.query(models.MockExam).filter(models.MockExam.id == exam_id).first()
    if not exam:
        raise HTTPException(status_code=404, detail="Mock exam not found")
    assert_owner(db, current_user, exam.subject_id)
    db.delete(exam)
    db.commit()
    return {"message": "Mock exam deleted successfully"}

@app.post("/api/subjects/{subject_id}/mock-exams", response_model=schemas.MockExamOut)
def generate_subject_mock_exam(subject_id: int, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    subject = own_subject(db, current_user, subject_id)

    materials = db.query(models.Material).filter(models.Material.subject_id == subject_id).all()
    texts = [m.summary for m in materials if m.summary]

    # Generate questions
    generated = mock_exam_agent.generate_mock_exam(subject.name, texts)
    questions_list = generated.get("questions", [])

    # Create exam in DB
    exam = models.MockExam(
        subject_id=subject_id,
        score=0.0,
        duration_seconds=0,
        status="in_progress"
    )
    db.add(exam)
    db.commit()
    db.refresh(exam)

    # Add questions
    for idx, q in enumerate(questions_list):
        db_q = models.MockExamQuestion(
            mock_exam_id=exam.id,
            question=q.get("question", "Conceptual essay question."),
            reference_source=q.get("reference_source", "Lecture Notes")
        )
        db.add(db_q)

    db.commit()
    db.refresh(exam)
    return exam

@app.post("/api/mock-exams/{exam_id}/submit", response_model=schemas.MockExamOut)
def submit_mock_exam(exam_id: int, req: schemas.MockExamSubmitRequest, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    exam = db.query(models.MockExam).filter(models.MockExam.id == exam_id).first()
    if not exam:
        raise HTTPException(status_code=404, detail="Mock exam session not found")
    assert_owner(db, current_user, exam.subject_id)

    materials = db.query(models.Material).filter(models.Material.subject_id == exam.subject_id).all()
    texts = [m.summary for m in materials if m.summary]

    # Map user answers
    answers_map = {ans.question_id: ans.user_answer for ans in req.answers}

    # Grade
    questions_payload = []
    user_answers_payload = []
    db_questions = exam.questions

    for db_q in db_questions:
        ans = answers_map.get(db_q.id, "")
        questions_payload.append({
            "id": db_q.id,
            "question": db_q.question,
            "reference_source": db_q.reference_source
        })
        user_answers_payload.append(ans)

    graded = mock_exam_agent.grade_mock_exam(questions_payload, user_answers_payload, texts)
    overall_score = graded.get("overall_score", 0.0)

    # Save to questions
    graded_map = {g.get("question_id"): g for g in graded.get("graded_questions", [])}
    for db_q in db_questions:
        ans = answers_map.get(db_q.id, "")
        g_data = graded_map.get(db_q.id, {})
        db_q.user_answer = ans
        db_q.ai_grade = g_data.get("ai_grade", 50.0)
        db_q.ai_feedback = g_data.get("ai_feedback", "Recall reviewed by coach.")
        db_q.reference_source = g_data.get("reference_source", db_q.reference_source)

    # Update exam
    exam.score = overall_score
    exam.duration_seconds = req.duration_seconds
    exam.status = "graded"
    exam.completed_at = datetime.datetime.utcnow().isoformat()

    # Slightly update overall subject readiness confidence
    subject = exam.subject
    if subject:
        subject.confidence_score = min(100.0, subject.confidence_score + (overall_score * 0.12))

    db.commit()
    db.refresh(exam)
    return exam

@app.get("/api/subjects/{subject_id}/formulas", response_model=List[schemas.FormulaOut])
def get_subject_formulas(subject_id: int, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    assert_owner(db, current_user, subject_id)
    # List-only: generation now happens explicitly via /formulas/generate
    return db.query(models.Formula).filter(models.Formula.subject_id == subject_id).all()

@app.post("/api/subjects/{subject_id}/formulas/generate", response_model=List[schemas.FormulaOut])
def generate_cheat_sheet(subject_id: int, req: schemas.GenerateCheatSheetRequest, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    """Generate cheat-sheet entries from the user's selected materials.
    Empty material_ids means 'use all materials for this subject'."""
    subject = own_subject(db, current_user, subject_id)

    materials_q = db.query(models.Material).filter(models.Material.subject_id == subject_id)
    if req.material_ids:
        materials_q = materials_q.filter(models.Material.id.in_(req.material_ids))
    materials = materials_q.all()
    if not materials:
        raise HTTPException(status_code=400, detail="No matching materials to generate from. Upload materials first.")

    # Prefer the richer deep-research summary when available
    texts = [(m.deep_research_summary or m.summary) for m in materials if (m.deep_research_summary or m.summary)]
    if not texts:
        raise HTTPException(status_code=400, detail="Selected materials have no processed content yet — wait for ingestion to finish.")

    extracted = formula_extractor_agent.extract_formulas(subject.name, texts)
    formulas_list = extracted.get("formulas", [])
    if not formulas_list:
        raise HTTPException(status_code=502, detail="The AI couldn't extract formulas from the selected materials. Try different resources.")

    if req.replace_existing:
        db.query(models.Formula).filter(models.Formula.subject_id == subject_id).delete()

    for f in formulas_list:
        db_f = models.Formula(
            subject_id=subject_id,
            name=f.get("name", "Key Theorem"),
            latex_code=f.get("latex_code", ""),
            description=f.get("description", ""),
            variables_json=json.dumps(f.get("variables", [])),
            derivation_steps_json=json.dumps(f.get("derivation_steps", []))
        )
        db.add(db_f)
    db.commit()

    return db.query(models.Formula).filter(models.Formula.subject_id == subject_id).all()

@app.post("/api/subjects/{subject_id}/formulas", response_model=schemas.FormulaOut)
def create_formula(subject_id: int, formula: schemas.FormulaCreate, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    assert_owner(db, current_user, subject_id)
    db_f = models.Formula(subject_id=subject_id, **formula.model_dump())
    db.add(db_f)
    db.commit()
    db.refresh(db_f)
    return db_f

@app.put("/api/formulas/{formula_id}", response_model=schemas.FormulaOut)
def update_formula(formula_id: int, updates: schemas.FormulaUpdate, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    formula = db.query(models.Formula).filter(models.Formula.id == formula_id).first()
    if not formula:
        raise HTTPException(status_code=404, detail="Formula entry not found")
    assert_owner(db, current_user, formula.subject_id)
    for key, value in updates.model_dump(exclude_unset=True).items():
        setattr(formula, key, value)
    db.commit()
    db.refresh(formula)
    return formula

@app.delete("/api/formulas/{formula_id}")
def delete_formula(formula_id: int, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    formula = db.query(models.Formula).filter(models.Formula.id == formula_id).first()
    if not formula:
        raise HTTPException(status_code=404, detail="Formula entry not found")
    assert_owner(db, current_user, formula.subject_id)
    db.delete(formula)
    db.commit()
    return {"message": "Formula deleted successfully"}

@app.post("/api/formulas/{formula_id}/note")
def update_formula_note(formula_id: int, note_data: dict, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    formula = db.query(models.Formula).filter(models.Formula.id == formula_id).first()
    if not formula:
        raise HTTPException(status_code=404, detail="Formula entry not found")
    assert_owner(db, current_user, formula.subject_id)

    note = note_data.get("note", "")
    # Append to description or keep track
    formula.description = (formula.description or "") + f"\n\n*Student Study Note:* {note}"
    db.commit()
    db.refresh(formula)
    return {"message": "Study note updated successfully!", "description": formula.description}

